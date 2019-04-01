"""
See http://pbpython.com/creating-powerpoint.html for details on this script
Requires https://python-pptx.readthedocs.org/en/latest/index.html
Example program showing how to read in Excel, process with pandas and
output to a PowerPoint file.
"""

from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches
import argparse
import numpy as np
from datetime import date
import jsonschema, json
import matplotlib.pyplot as plt
from summarize_metrics import summarize_data, BAAThresholds
from pathlib import Path
from argparse import ArgumentParser, ArgumentError
import pandas as pd


def directory_type(arg_string):
    """
    Allows arg parser to handle directories
    :param arg_string: A path, relative or absolute to a folder
    :return: A python pure path object to a directory.
    """
    directory_path = Path(arg_string)
    if directory_path.exists() and directory_path.is_dir():
        return str(directory_path)
    raise ArgumentError("{} is not a valid directory.".format(arg_string))


def list_type(arg_string):
    """
    Allows arg parser to handle strings
    :param arg_string: a python syntaxed list, as a string
    :return: A python list
    """
    arguments = arg_string.split()
    return arguments


def parse_args():
    """ Setup the input and output arguments for the script
    Return the parsed input and output files
    """
    parser = argparse.ArgumentParser(description='Create ppt report')
    parser.add_argument('-infile',
                        type=argparse.FileType('r'),
                        help='Powerpoint file used as the template',
                        required=True)
    parser.add_argument('-outfile',
                        type=argparse.FileType('w'),
                        help='Output powerpoint report file',
                        required=True)
    parser.add_argument('-rootdir',
                        type=directory_type,
                        help='Root directory for all teams',
                        required=True)
    parser.add_argument('-teams',
                        nargs='+',
                        help='Name of all teams meant to be evaluated',
                        required=True)
    parser.add_argument('-aois',
                        nargs='+',
                        help='All AOIs meant to be evaluated',
                        required=True)
    return parser.parse_args()


def create_ppt(input, output, summarized_results, baa_thresolds):
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    prs = Presentation(input)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "CORE3D Metrics Report"
    subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())
    subbox = slide.placeholders[13]
    subbox.text = "JHU/APL"

    # Create Table of all teams
    # Create columns from summarized results
    metrics_names = {'Metrics': ['2D Correctness', '2D Completeness', '2D IOU', '3D Correctness',
                     '3D Completeness', '3D IOU']}
    phase_1a_thresholds = {'Phase 1A Thresholds': [baa_thresolds.correctness_2d[0],
                                                   baa_thresolds.completeness_2d[0],
                                                   baa_thresolds.jaccard_index_2d[0],
                                                   baa_thresolds.correctness_3d[0],
                                                   baa_thresolds.completeness_3d[0],
                                                   baa_thresolds.jaccard_index_3d[0]]}
    phase_2b_thresholds = {'Phase 1A Thresholds': [baa_thresolds.correctness_2d[3],
                                                   baa_thresolds.completeness_2d[3],
                                                   baa_thresolds.jaccard_index_2d[3],
                                                   baa_thresolds.correctness_3d[3],
                                                   baa_thresolds.completeness_3d[3],
                                                   baa_thresolds.jaccard_index_3d[3]]}
    df_metrics_names = pd.DataFrame(data=metrics_names)
    df_1a_thresholds = pd.DataFrame(data=phase_1a_thresholds)
    df_2b_thresholds = pd.DataFrame(data=phase_2b_thresholds)
    metrics_column = {}
    for team in summarized_results:
        metrics_2D = [summarized_results[team][6]['2D']["correctness"],
                      summarized_results[team][6]['2D']["completeness"],
                      summarized_results[team][6]['2D']["jaccardindex"]]
        metrics_3D = [summarized_results[team][6]['3D']["correctness"],
                      summarized_results[team][6]['3D']["completeness"],
                      summarized_results[team][6]['3D']["jaccardindex"]]
        metrics_column = {team: metrics_2D+metrics_3D}
    df_performer_metrics = pd.DataFrame(data=metrics_column)
    df_mean_scores = pd.concat([df_metrics_names,
                                df_1a_thresholds,
                                df_2b_thresholds,
                                df_performer_metrics],
                               axis=1)

    for team in summarized_results:
        summary_metrics_slide_layout = prs.slide_layouts[4]
        slide = prs.slides.add_slide(summary_metrics_slide_layout)
        title = slide.shapes.title
        current_date = date.today()
        title.text = "Mean Scores - {0}, Self Test".format(current_date.strftime('%B'))
    # Add Metrics Table
    for team in summarized_results:
        summary_metrics_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(summary_metrics_slide_layout)
        title = slide.shapes.title
        title.text = "Summary of Results for " + team

    prs.save(output)


def main():
    args = parse_args()
    baa_thresholds = BAAThresholds()
    summarized_results = summarize_data(baa_thresholds, args.rootdir, args.teams, args.aois)
    create_ppt(args.infile.name, args.outfile.name, summarized_results, baa_thresholds)

if __name__ == "__main__":
    main()
